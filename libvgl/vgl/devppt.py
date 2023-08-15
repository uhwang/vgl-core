'''
    devpptx.py
    
    7/23/2023
    
    Line
    
    Use Case 1 - Specify all points in slide coordinates
    A very simple case can be had by specifying the verticies in slide coordinates:
    
    >>> freeform_builder = shapes.build_freeform(Inches(1.5), Inches(1))
    >>> freeform_builder.add_line_segments((
    ...     (Inches(2),   Inches(2)),
    ...     (Inches(1),   Inches(2)),
    ...     (Inches(1.5), Inches(1)),
    ... ))
    >>> freeform_shape = freeform_builder.convert_to_shape()
    >>> freeform_shape.left.inches, freeform_shape.top.inches
    1.0, 1.0
    >>> freeform_shape.width.inches, freeform_shape.height.inches
    1.0, 1.0
    
    Polygon
    
    >>> vertices = ((100, 200), (200, 300), (300, 400))
    >>> freeform_builder.add_line_segments(vertices, close=True)

    
    chat.openai.com
    ================
    
    line_shape = slide.shapes.add_shape(9, left, top, width, height)
    line_shape.shadow.inherit = False
    line_shape.shadow.blur_radius = 0
    line_shape.shadow.distance = 0
    
    after creating the line shape using slide.shapes.add_shape(), we access the shadow property of the shape and set the inherit, blur_radius, and distance attributes to disable the shadow effect. By setting inherit to False, we make sure that the shape's shadow properties are independent and not inherited from any style.
    
    With these changes, the line will be drawn without any shadow effect in the resulting PowerPoint presentation.
    
    
    chat.openai.com
    ===============
    
    insert shapes to a group shape
    group_shape = slide.shapes.add_group_shape()
    
    shapes = [
        slide.shapes.add_shape(
            pptx.enum.shapes.MSO_SHAPE.RECTANGLE,
            left=100,
            top=100,
            width=100,
            height=50),
        slide.shapes.add_shape(
            pptx.enum.shapes.MSO_SHAPE.RECTANGLE,
            left=200,
            top=100,
            width=100,
            height=150)
    ]
    
    for shape in shapes:
        group_shape.shapes._spTree.insert_element_before(shape._element, 'p:extLst')
    
'''
import collections 
import collections.abc
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import numpy as np

from . import color
from . import device
from . import drvemf 
from . import linepat
from . import patline
from . import gdiobj

def Polyline(slide, x, y, lcol, lthk, lpat, fcol, closed):

    free_form = slide.shapes.build_freeform(Inches(x[0]), Inches(y[0]))
    free_form.add_line_segments(
        [(Inches(x1), Inches(y1)) for x1, y1 in zip(x[1:],y[1:])],
        close=closed
    )
    
    line_shape = free_form.convert_to_shape()
    
    if closed == False: 
        line_shape.fill.background()
    else:
        if isinstance(fcol, color.Color):
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = RGBColor(fcol.r, fcol.g, fcol.b)
        else:
            line_shape.fill.background()    
    
    if lcol:
        if fcol and lcol == fcol:
            line_shape.line.color.rgb = RGBColor(fcol.r, fcol.g, fcol.b)
            line_shape.line.width = Inches(lthk)
        else:
            line_shape.line.color.rgb = RGBColor(lcol.r, lcol.g, lcol.b)
            line_shape.line.width = Inches(lthk)
        
    line_shape.shadow.inherit     = False
    line_shape.shadow.blur_radius = 0
    line_shape.shadow.distance    = 0
    
# Not working. 

#def Line(slide, x1, y1, x2, y2, lcol, lthk, lpat):
#    line_shape= slide.shapes.add_shape(_SHAPE_LINE, 
#                                        Inches(x1), 
#                                        Inches(y1), 
#                                        Inches(x2-x1), 
#                                        Inches(y2-y1))
#    line_shape.line.color.rgb = RGBColor(lcol.r, lcol.g, lcol.b)
#    line_shape.line.width = Pt(lthk)
#    line_shape.fill.background()
#    line_shape.shadow.inherit = False
#    line_shape.shadow.blur_radius = 0
#    line_shape.shadow.distance = 0
#    line_shape.rotation = -45000  # Rotation angle is in 1/60000 of a degree
    
    
class DevicePPT(device.DeviceVector):
    def __init__(self, fname, gbox):
        super().__init__()
        self.fname = fname
        self.ppt = pptx.Presentation()
        self.blank_slide_layout = self.ppt.slide_layouts[6]
        self.slide = self.ppt.slides.add_slide(self.blank_slide_layout)
        self.pen = gdiobj.Pen()
        self.brush = gdiobj.Brush()
        
    def set_device(self, frm, extend=device._FIT_NONE):
        self.frm = frm
        self.set_plot(frm,extend)
        
    def close(self):
        self.ppt.save(self.fname)
            
    def make_pen(self, lcol, lthk):
        self.pen.lcol = lcol
        self.pen.lthk = lthk
    
    def delete_pen(self):
        self.pen.lcol = None
        self.pen.lthk = None
        
    def circle(self, x,y, rad, lcol=color.BLACK, lthk=0.001, lpat=linepat._PAT_SOLID, fcol=None):
        rrad = np.linspace(0, np.pi*2, self._circle_point)
        x1 = x+rad*np.cos(rrad)
        y1 = y+rad*np.sin(rrad)
        self._polyline(x1,y1, lcol, lthk*self.frm.hgt(), lpat, fcol, closed=True, viewport=False)
        
    def lcircle(self, x,y, rad, lcol=color.BLACK, lthk=0.001, lpat=linepat._PAT_SOLID, fcol=None):
        rrad = np.linspace(0, np.pi*2, self._circle_point)
        x1 = x+rad*np.cos(rrad)
        y1 = y+rad*np.sin(rrad)
        self._polyline(x,y, rad, lcol, lthk*self.frm.hgt(), lpat, fcol, closed=True, viewport=True)
        
    def line(self, sx, sy, ex, ey, lcol=color.BLACK, lthk=0.001, lpat=linepat._PAT_SOLID, viewport=False):
        if self.pen.lcol:
            _lcol = self.pen.lcol
            _lthk = self.pen.lthk
        else:
            _lcol = lcol
            _lthk = lthk*self.frm.hgt()
        self._polyline((sx,ex), (sy,ey), _lcol, _lthk, lpat, None, False, viewport=False)
        
    def lline(self, sx, sy, ex, ey, lcol=color.BLACK, lthk=0.001, lpat=linepat._PAT_SOLID, viewport=True):
        if self.pen.lcol:
            _lcol = self.pen.lcol
            _lthk = self.pen.lthk
        else:
            _lcol = lcol
            _lthk = lthk*self.frm.hgt()
        #print(_lthk)
        self._polyline((sx,ex), (sy,ey), _lcol, _lthk, lpat, None, False, viewport=True)
        
    def _polyline(self, x, y, 
                        lcol=None, 
                        lthk=None, 
                        lpat=linepat._PAT_SOLID, 
                        fcol=None, 
                        closed=False, 
                        viewport=False):
        pat_inst = isinstance(lpat, linepat.LinePattern)

        #if lthk: _lthk = lthk*self.frm.hgt()
        #else: _lthk = 0

        if viewport:
            px, py = x, y
        else:
            px = [self._x_viewport(xx) for xx in x]
            py = [self._y_viewport(yy) for yy in y]

        # polygon and solid outline
        if fcol and closed and pat_inst==linepat._PAT_SOLID:
            Polyline(self.slide, px, py, lcol, lthk, lpat, fcol, closed)
            
        # polyline/polygon and solid/patterened outline
        else:
            if pat_inst:
                # fill the polygon
                if fcol and closed:
                    lcol1 = fcol
                    Polyline(self.slide, px, py, lcol1, lthk, lpat, fcol, closed)
                    
                # polygon & outline (one more point to close)
                if closed:
                    if isinstance(x, np.ndarray):
                        xp = np.append(x, x[0])
                        yp = np.append(y, y[0])
                    elif isinstance(x, list):
                        xp = x.copy()
                        yp = y.copy()
                        xp.append(x[0])
                        yp.append(y[0])
                else:
                    xp, yp = x, y
                    
                if viewport:
                    pat_seg = patline.get_pattern_line(self, xp, yp, lpat.pat_len, lpat.pat_t, viewport=True)
                else:
                    pat_seg = patline.get_pattern_line(self, xp, yp, lpat.pat_len, lpat.pat_t)
    
                for p1 in pat_seg:
                    x2 = [p2[0] for p2 in p1 ]
                    y2 = [p2[1] for p2 in p1 ]
                    Polyline(self.slide, x2, y2, lcol, lthk, lpat, fcol, closed=False)
            else:
                Polyline(self.slide, px, py, lcol, lthk, lpat, fcol, closed)
         
    def polyline(self, x, y, lcol=color.BLACK, lthk=0.001, lpat=linepat._PAT_SOLID, closed=False):
        if self.pen.lcol:
            _lcol = self.pen.lcol
            _lthk = self.pen.lthk
        else:
            _lcol = lcol
            _lthk = lthk*self.frm.hgt()
        self._polyline(x, y, lcol, lthk, lpat, None, closed, False)
        
    def lpolyline(self, x, y, lcol=color.BLACK, lthk=0.001, lpat=linepat._PAT_SOLID, closed=False):
        if self.pen.lcol:
            _lcol = self.pen.lcol
            _lthk = self.pen.lthk
        else:
            _lcol = lcol
            _lthk = lthk*self.frm.hgt()
        self._polyline(x, y, lcol, lthk, lpat, None, closed, True)
    
    def polygon(self, x, y, lcol=color.BLACK, lthk=0.001, lpat=linepat._PAT_SOLID, fcol=None):
        if self.pen.lcol:
            _lcol = self.pen.lcol
            _lthk = self.pen.lthk
        else:
            _lcol = lcol
            _lthk = lthk*self.frm.hgt()    
        self._polyline(x, y, _lcol, _lthk, lpat, fcol, True, viewport=False)
        
    def lpolygon(self, x, y, lcol=color.BLACK, lthk=0.001, lpat=linepat._PAT_SOLID, fcol=None):
        if self.pen.lcol:
            _lcol = self.pen.lcol
            _lthk = self.pen.lthk
        else:
            _lcol = lcol
            _lthk = lthk*self.frm.hgt()    
        self._polyline(x, y, _lcol, _lthk, lpat, fcol, True, viewport=True)
        
    def begin_symbol(self, sym):
        pass
        
    def end_symbol(self):
        pass        
        
    def symbol(self, x,y,sym,draw=False):
        px, py = sym.update_xy( self._x_viewport(x),
                                self._y_viewport(y) )
        self._polyline(px, py, sym.lcol, sym.lthk*self.frm.hgt(), 
        linepat._PAT_SOLID, sym.fcol, closed=True, viewport=True)
        
    def stroke(self):
        pass